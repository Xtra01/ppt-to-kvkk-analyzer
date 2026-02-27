"""
PPT → Vektör Dönüştürücü
========================
Kaynaklar klasöründeki PowerPoint dosyalarını AI'ın anlayabileceği
semantik vektörlere (embeddings) dönüştürür.

Kullanım:
    python ppt_to_vectors.py --all                   # Tam işlem: çıkar + vektörleştir
    python ppt_to_vectors.py --extract               # Sadece metin çıkar
    python ppt_to_vectors.py --vectorize             # Sadece vektörleştir
    python ppt_to_vectors.py --txt                   # TXT dosyalarına aktar
    python ppt_to_vectors.py --search "KVKK nedir?"  # Semantik arama yap
"""

__version__ = "1.2.0"

import os
import sys
import io
import json
import argparse
import logging
import time
from pathlib import Path
from typing import List, Dict, Tuple, Optional

# Windows konsolunda UTF-8 karakterleri düzgün yazdır
if sys.stdout.encoding and sys.stdout.encoding.lower() not in ("utf-8", "utf8"):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

# ── LOGGING AYARI ────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger(__name__)

# ── DİZİN AYARLARI ──────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent.parent
KAYNAKLAR_DIR    = BASE_DIR / "input"                 # PPT kaynak dosyaları
CIKTILAR_DIR     = BASE_DIR / "output" / "vectors"    # vektör çıktıları
TXT_CIKTILAR_DIR = BASE_DIR / "output" / "txt"        # düz metin çıktıları

# ── MODEL AYARLARI ───────────────────────────────────────────────────────────
# Türkçe dahil 50+ dil destekleyen hafif ama etkili model
MODEL_NAME = "paraphrase-multilingual-MiniLM-L12-v2"

# ── CHUNK AYARLARI ───────────────────────────────────────────────────────────
CHUNK_MAX_CHARS = 500       # Bir parçanın maksimum karakter uzunluğu
MIN_CHUNK_CHARS = 20        # Bu uzunluktan kısa parçalar atlanır


# ═══════════════════════════════════════════════════════════════════════════════
# 1) PPTX METİN ÇIKARMA
# ═══════════════════════════════════════════════════════════════════════════════

def _shape_texts(shape) -> List[str]:
    """Bir shape'ten tüm metinleri çıkarır (text frame + tablo)."""
    texts: List[str] = []

    # Normal metin çerçevesi
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            line = paragraph.text.strip()
            if line:
                texts.append(line)

    # Tablo hücreleri
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                texts.append(" | ".join(row_texts))

    # Grup şekilleri – iç içe shape'lere eriş
    if shape.shape_type is not None:
        try:
            if hasattr(shape, "shapes"):
                for child in shape.shapes:
                    texts.extend(_shape_texts(child))
        except Exception:
            pass

    return texts


def extract_pptx(file_path: Path) -> List[Dict]:
    """Bir PPTX dosyasından slayt bazında metin çıkarır."""
    from pptx import Presentation

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        logger.error(f"  ✗ Dosya okunamadı: {file_path.name} → {exc}")
        return []

    slides: List[Dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        all_texts: List[str] = []
        for shape in slide.shapes:
            all_texts.extend(_shape_texts(shape))

        full_text = "\n".join(all_texts)
        if full_text.strip():
            slides.append({
                "dosya": file_path.name,
                "slayt_no": slide_idx,
                "metin": full_text.strip(),
            })
    return slides


def extract_all(kaynak_dir: Path) -> List[Dict]:
    """Klasördeki tüm PPTX dosyalarından metin çıkarır."""
    pptx_files = sorted(kaynak_dir.glob("*.pptx"))
    if not pptx_files:
        logger.warning(f"PPTX dosyası bulunamadı: {kaynak_dir}")
        return []

    logger.info(f"{len(pptx_files)} PPTX dosyası bulundu.")
    all_slides: List[Dict] = []

    for f in pptx_files:
        logger.info(f"  → {f.name}")
        slides = extract_pptx(f)
        all_slides.extend(slides)
        logger.info(f"    {len(slides)} slayt çıkarıldı")

    logger.info(f"Toplam {len(all_slides)} slayt metni elde edildi.")
    return all_slides


# ═══════════════════════════════════════════════════════════════════════════════
# 2) METİN PARÇALAMA  (Chunking)
# ═══════════════════════════════════════════════════════════════════════════════

def _chunk_text(text: str, max_chars: int = CHUNK_MAX_CHARS) -> List[str]:
    """
    Metni paragraf sınırlarına göre parçalara böler.
    Paragraf tek başına sığmıyorsa kelime bazlı keser.
    """
    if len(text) <= max_chars:
        return [text]

    chunks: List[str] = []
    paragraphs = text.split("\n")
    buf = ""

    for para in paragraphs:
        candidate = (buf + "\n" + para).strip() if buf else para
        if len(candidate) <= max_chars:
            buf = candidate
        else:
            if buf:
                chunks.append(buf)
            # Paragraf tek başına çok uzunsa kelime bazlı böl
            if len(para) > max_chars:
                words = para.split()
                buf = ""
                for w in words:
                    test = (buf + " " + w).strip() if buf else w
                    if len(test) <= max_chars:
                        buf = test
                    else:
                        if buf:
                            chunks.append(buf)
                        buf = w
            else:
                buf = para

    if buf.strip():
        chunks.append(buf.strip())

    return chunks


def create_chunks(slides: List[Dict]) -> List[Dict]:
    """Slayt metinlerinden indekslenmiş parçalar oluşturur."""
    chunks: List[Dict] = []
    cid = 0

    for slide in slides:
        parts = _chunk_text(slide["metin"])
        for part_no, text in enumerate(parts, start=1):
            if len(text) < MIN_CHUNK_CHARS:
                continue
            chunks.append({
                "id": cid,
                "dosya": slide["dosya"],
                "slayt_no": slide["slayt_no"],
                "parca_no": part_no,
                "metin": text,
            })
            cid += 1

    return chunks


# ═══════════════════════════════════════════════════════════════════════════════
# 3) VEKTÖRLEŞTİRME  (Embedding)
# ═══════════════════════════════════════════════════════════════════════════════

def vectorize(chunks: List[Dict], model_name: str = MODEL_NAME):
    """
    Chunk metinlerini semantik vektörlere dönüştürür.
    numpy ndarray (N, dim) döner.
    """
    import numpy as np
    from sentence_transformers import SentenceTransformer

    logger.info(f"Model yükleniyor: {model_name}")
    t0 = time.time()
    model = SentenceTransformer(model_name)
    logger.info(f"Model hazır ({time.time() - t0:.1f}s)")

    texts = [c["metin"] for c in chunks]
    logger.info(f"{len(texts)} parça vektörleştirilecek...")

    embeddings = model.encode(
        texts,
        show_progress_bar=True,
        batch_size=32,
        normalize_embeddings=True,   # cosine similarity için normalize
    )

    logger.info(f"Vektör matrisi: {embeddings.shape}  (parça × boyut)")
    return embeddings


# ═══════════════════════════════════════════════════════════════════════════════
# 4) KAYDETME
# ═══════════════════════════════════════════════════════════════════════════════

def save(embeddings, chunks: List[Dict], out_dir: Path):
    """Vektörleri (.npy) ve metadata'yı (.json) diske kaydeder."""
    import numpy as np

    out_dir.mkdir(parents=True, exist_ok=True)

    # Vektörler
    vec_path = out_dir / "vectors.npy"
    np.save(str(vec_path), embeddings)
    logger.info(f"Vektörler → {vec_path}")

    # Metadata + chunk'lar
    meta = {
        "model": MODEL_NAME,
        "vektor_boyutu": int(embeddings.shape[1]),
        "toplam_parca": len(chunks),
        "chunks": chunks,
    }
    meta_path = out_dir / "metadata.json"
    with open(str(meta_path), "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    logger.info(f"Metadata → {meta_path}")

    # İnsan tarafından okunabilir özet
    _write_summary(embeddings, chunks, out_dir)


def _write_summary(embeddings, chunks: List[Dict], out_dir: Path):
    """Özet rapor dosyası oluşturur."""
    dosyalar = sorted(set(c["dosya"] for c in chunks))
    summary_path = out_dir / "ozet_rapor.txt"

    with open(str(summary_path), "w", encoding="utf-8") as f:
        f.write("PPT → VEKTÖR DÖNÜŞTÜRME RAPORU\n")
        f.write("=" * 50 + "\n\n")
        f.write(f"Model         : {MODEL_NAME}\n")
        f.write(f"Vektör boyutu : {embeddings.shape[1]}\n")
        f.write(f"Toplam parça  : {len(chunks)}\n")
        f.write(f"Dosya sayısı  : {len(dosyalar)}\n\n")

        for dosya in dosyalar:
            n = sum(1 for c in chunks if c["dosya"] == dosya)
            f.write(f"  • {dosya} → {n} parça\n")

        f.write("\n" + "=" * 50 + "\n")
        f.write("Çıktı dosyaları:\n")
        f.write("  vectors.npy   – numpy vektör matrisi\n")
        f.write("  metadata.json – metin parçaları ve metadata\n")
        f.write("  ozet_rapor.txt – bu rapor\n")

    logger.info(f"Özet rapor → {summary_path}")


# ═══════════════════════════════════════════════════════════════════════════════
# 5) SEMANTİK ARAMA
# ═══════════════════════════════════════════════════════════════════════════════

def search(query: str, out_dir: Path, top_k: int = 5, model_name: str = MODEL_NAME):
    """Vektör deposunda cosine similarity ile semantik arama yapar."""
    import numpy as np
    from sentence_transformers import SentenceTransformer

    vec_path = out_dir / "vectors.npy"
    meta_path = out_dir / "metadata.json"

    if not vec_path.exists() or not meta_path.exists():
        logger.error("Vektör dosyaları bulunamadı. Önce --all ile dönüştürme yapın.")
        sys.exit(1)

    embeddings = np.load(str(vec_path))
    with open(str(meta_path), "r", encoding="utf-8") as f:
        metadata = json.load(f)
    chunks = metadata["chunks"]

    model = SentenceTransformer(model_name)
    q_vec = model.encode([query], normalize_embeddings=True)

    # Cosine similarity  (normalize vektörler → dot product)
    scores = (embeddings @ q_vec.T).flatten()
    top_idx = scores.argsort()[::-1][:top_k]

    print()
    print("=" * 60)
    print(f"  ARAMA: \"{query}\"")
    print("=" * 60)

    for rank, idx in enumerate(top_idx, start=1):
        c = chunks[idx]
        score = scores[idx]
        print(f"\n  #{rank}  Benzerlik: {score:.4f}")
        print(f"  Dosya : {c['dosya']}  |  Slayt: {c['slayt_no']}")
        print(f"  ─────────────────────────────────────────────")
        # Metni 300 karakterle sınırla
        preview = c["metin"][:300]
        if len(c["metin"]) > 300:
            preview += " …"
        print(f"  {preview}")

    print("\n" + "=" * 60 + "\n")


# ═══════════════════════════════════════════════════════════════════════════════
# 6) TXT DIŞA AKTARMA
# ═══════════════════════════════════════════════════════════════════════════════

def export_txt(kaynak_dir: Path, txt_dir: Path):
    """Her PPTX dosyasını ayrı bir .txt dosyasına aktarır."""
    from pptx import Presentation

    pptx_files = sorted(kaynak_dir.glob("*.pptx"))
    if not pptx_files:
        logger.warning(f"PPTX dosyası bulunamadı: {kaynak_dir}")
        return

    txt_dir.mkdir(parents=True, exist_ok=True)
    logger.info(f"{len(pptx_files)} PPTX dosyası → TXT dönüştürülecek")

    for f in pptx_files:
        logger.info(f"  → {f.name}")
        try:
            prs = Presentation(str(f))
        except Exception as exc:
            logger.error(f"  ✗ Dosya okunamadı: {f.name} → {exc}")
            continue

        lines: List[str] = []
        lines.append(f"{'=' * 60}")
        lines.append(f"  {f.stem}")
        lines.append(f"{'=' * 60}")
        lines.append("")

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts: List[str] = []
            for shape in slide.shapes:
                slide_texts.extend(_shape_texts(shape))

            if slide_texts:
                lines.append(f"── Slayt {slide_idx} {'─' * 40}")
                for t in slide_texts:
                    lines.append(t)
                lines.append("")

        # Dosyayı yaz
        txt_name = f.stem + ".txt"
        txt_path = txt_dir / txt_name
        with open(str(txt_path), "w", encoding="utf-8") as out:
            out.write("\n".join(lines))
        logger.info(f"    → {txt_path.name} ({len(prs.slides)} slayt)")

    logger.info(f"TXT çıktıları → {txt_dir}")


# ═══════════════════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="PPT dosyalarını AI vektörlerine dönüştürme aracı",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Örnekler:
  python ppt_to_vectors.py --all                   Tam işlem
  python ppt_to_vectors.py --search "kişisel veri"  Semantik arama
  python ppt_to_vectors.py --search "KVKK yaptırımlar" --top-k 10
        """,
    )
    parser.add_argument("--extract",   action="store_true", help="PPT → metin çıkarma")
    parser.add_argument("--vectorize", action="store_true", help="Metinler → vektör")
    parser.add_argument("--all",       action="store_true", help="Çıkarma + vektörleştirme")
    parser.add_argument("--txt",       action="store_true", help="PPT → TXT dosyalarına aktar")
    parser.add_argument("--search",    type=str,            help="Semantik arama sorgusu")
    parser.add_argument("--top-k",     type=int, default=5, help="Arama sonuç sayısı")

    args = parser.parse_args()

    if not any([args.extract, args.vectorize, args.all, args.txt, args.search]):
        parser.print_help()
        return

    chunks = None

    # ── ADIM 1: Metin çıkarma ────────────────────────────────────────────
    if args.extract or args.all:
        logger.info("━" * 50)
        logger.info("ADIM 1 · PPT'lerden metin çıkarılıyor")
        logger.info("━" * 50)

        slides = extract_all(KAYNAKLAR_DIR)
        if not slides:
            logger.error("Hiç metin çıkarılamadı, işlem durduruluyor.")
            sys.exit(1)

        chunks = create_chunks(slides)
        logger.info(f"{len(chunks)} metin parçası oluşturuldu.")

        # Ara çıktı – sonraki adımda lazım olabilir
        CIKTILAR_DIR.mkdir(parents=True, exist_ok=True)
        ext_path = CIKTILAR_DIR / "extracted_chunks.json"
        with open(str(ext_path), "w", encoding="utf-8") as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)
        logger.info(f"Çıkarılan metinler → {ext_path}")

    # ── ADIM 2: Vektörleştirme ───────────────────────────────────────────
    if args.vectorize or args.all:
        logger.info("━" * 50)
        logger.info("ADIM 2 · Vektörleştirme")
        logger.info("━" * 50)

        # Chunk'lar bellekte yoksa diskten oku
        if chunks is None:
            ext_path = CIKTILAR_DIR / "extracted_chunks.json"
            if not ext_path.exists():
                logger.error("Çıkarılmış metin bulunamadı. Önce --extract çalıştırın.")
                sys.exit(1)
            with open(str(ext_path), "r", encoding="utf-8") as f:
                chunks = json.load(f)

        embeddings = vectorize(chunks)
        save(embeddings, chunks, CIKTILAR_DIR)

        logger.info("━" * 50)
        logger.info("İşlem tamamlandı!")
        logger.info(f"Çıktılar: {CIKTILAR_DIR}")
        logger.info("━" * 50)

    # ── TXT DIŞA AKTARMA ─────────────────────────────────────────────────
    if args.txt:
        logger.info("━" * 50)
        logger.info("TXT · PPT'ler metin dosyalarına aktarılıyor")
        logger.info("━" * 50)
        export_txt(KAYNAKLAR_DIR, TXT_CIKTILAR_DIR)
        logger.info("━" * 50)
        logger.info("TXT dışa aktarma tamamlandı!")
        logger.info("━" * 50)

    # ── ARAMA ─────────────────────────────────────────────────────────────
    if args.search:
        search(args.search, CIKTILAR_DIR, top_k=args.top_k)


if __name__ == "__main__":
    main()
